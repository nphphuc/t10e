#!/usr/bin/env python3
"""Extract text content from all PPTX files in the Slide directory."""

import os
from pptx import Presentation

SLIDE_DIR = "Slide"
OUTPUT_DIR = "slide_texts"
os.makedirs(OUTPUT_DIR, exist_ok=True)

pptx_files = sorted([
    f for f in os.listdir(SLIDE_DIR)
    if f.endswith(".pptx") and not f.startswith("~")
])

print(f"Found {len(pptx_files)} PPTX files to process.\n")

for fname in pptx_files:
    filepath = os.path.join(SLIDE_DIR, fname)
    safe_name = fname.replace(".pptx", "").replace(" ", "_").replace("&", "and")
    output_path = os.path.join(OUTPUT_DIR, f"{safe_name}.txt")
    
    try:
        prs = Presentation(filepath)
        lines = []
        lines.append(f"=== FILE: {fname} ===\n")
        lines.append(f"Total slides: {len(prs.slides)}\n")
        
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
            lines.append(f"Layout: {layout_name}")
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            row_texts.append(cell.text.strip())
                        lines.append(" | ".join(row_texts))
        
        content = "\n".join(lines)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(content)
        
        slide_count = len(prs.slides)
        char_count = len(content)
        print(f"[OK] {fname}: {slide_count} slides, {char_count} chars -> {output_path}")
    
    except Exception as e:
        print(f"[ERR] {fname}: {e}")

print("\nDone! All texts extracted.")
