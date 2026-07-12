from pptx import Presentation
import os
import json

slides_dir = "Slide"
output = []

pptx_files = sorted([f for f in os.listdir(slides_dir) if f.endswith('.pptx')])

for filename in pptx_files:
    filepath = os.path.join(slides_dir, filename)
    try:
        prs = Presentation(filepath)
        
        course = {
            "filename": filename,
            "title": filename.replace(".pptx", "").replace("_", " "),
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": i,
                "texts": []
            }
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_data["texts"].append(text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            slide_data["texts"].append(" | ".join(row_texts))
            
            if slide_data["texts"]:
                course["slides"].append(slide_data)
        
        output.append(course)
        print("[OK] " + filename + " (" + str(len(course['slides'])) + " slides)")
    
    except Exception as e:
        print("[ERR] " + filename + ": " + str(e))

with open("slide_content.json", "w", encoding="utf-8") as f:
    json.dump(output, f, ensure_ascii=False, indent=2)

print("\n[DONE] Saved to slide_content.json")
print("Total files: " + str(len(output)))
